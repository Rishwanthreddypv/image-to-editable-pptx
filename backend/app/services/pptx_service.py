from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
import os
from typing import List
from app.schemas.layer import Layer
from app.core.config import logger

class PPTXService:
    def create_presentation(self, layers: List[Layer], width: int = 1280, height: int = 720, background_color: str = "#ffffff") -> io.BytesIO:
        """
        Convert document layers into a PowerPoint presentation.
        """
        logger.info(f"Creating PPTX presentation from layers with background {background_color}")
        prs = Presentation()
        
        # Set slide dimensions for 16:9 (Standard in modern PPTX)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        scale_x = prs.slide_width / width
        scale_y = prs.slide_height / height

        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Apply background fill
        fill = slide.background.fill
        fill.solid()
        try:
            color_hex = background_color.lstrip('#')
            fill.fore_color.rgb = RGBColor.from_string(color_hex)
        except:
            fill.fore_color.rgb = RGBColor(255, 255, 255) # Fallback to White

        for layer in layers:
            try:
                if layer.type == 'text':
                    self._add_text(slide, layer, scale_x, scale_y)
                elif layer.type == 'table':
                    self._add_table(slide, layer, scale_x, scale_y)
                elif layer.type == 'image':
                    self._add_image(slide, layer, scale_x, scale_y)
                else:
                    self._add_shape(slide, layer, scale_x, scale_y)
            except Exception as e:
                logger.error(f"Failed to add layer {layer.id} to PPTX: {str(e)}")

        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    def _add_text(self, slide, layer, sx, sy):
        # Add a small width buffer (5%) to account for font rendering differences
        width_buffer = 1.05
        
        left = int(layer.geometry.x * sx)
        top = int(layer.geometry.y * sy)
        width = int(layer.geometry.w * sx * width_buffer)
        height = int(layer.geometry.h * sy)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Match frontend "pixel-perfect" look by removing default margins
        tf.margin_bottom = 0
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.word_wrap = True
        
        # Ensure text starts from the top
        tf.vertical_anchor = MSO_ANCHOR.TOP

        tf.text = layer.content.get('text', '')
        
        # Scaling factor: 1280px wide canvas maps to 960pt wide slide (13.333")
        # Therefore, 1px in editor = 0.75pt in PPTX
        PT_PER_PX = 0.75
        default_font_size_px = 16
        
        # Apply styles
        style = layer.style or {}
        p = tf.paragraphs[0]
        
        # Add spacing to prevent "clumsy" overlapping look
        tf.line_spacing = 1.1 # Tighten slightly to match web rendering
        p.space_after = Pt(2) # Minimal paragraph spacing
        
        font_size_px = style.get('fontSize', default_font_size_px)
        # Apply the 0.75 scale factor to the font size
        p.font.size = Pt(font_size_px * PT_PER_PX)
        p.font.name = style.get('fontFamily', 'Arial')
        p.font.bold = style.get('fontWeight') == 'bold'
        
        if 'color' in style:
            try:
                color_hex = style['color'].lstrip('#')
                p.font.color.rgb = RGBColor.from_string(color_hex)
            except: pass
            
        if 'align' in style:
            mapping = {
                'left': PP_ALIGN.LEFT, 
                'center': PP_ALIGN.CENTER, 
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            p.alignment = mapping.get(style['align'], PP_ALIGN.LEFT)
        else:
            p.alignment = PP_ALIGN.LEFT

    def _add_table(self, slide, layer, sx, sy):
        left = int(layer.geometry.x * sx)
        top = int(layer.geometry.y * sy)
        width = int(layer.geometry.w * sx)
        height = int(layer.geometry.h * sy)

        cells = layer.content.get('cells', [])
        if not cells: return

        rows = max(c.get('rowIndex', 0) for c in cells) + 1
        cols = max(c.get('colIndex', 0) for c in cells) + 1

        # add_table returns a GraphicFrame, the table object is in its .table property
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table

        for cell_data in cells:
            r, c = cell_data.get('rowIndex', 0), cell_data.get('colIndex', 0)
            if r < rows and c < cols:
                cell = table.cell(r, c)
                cell.text = str(cell_data.get('content', ''))
                # Style the cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.name = 'Arial'

    def _add_image(self, slide, layer, sx, sy):
        from PIL import Image
        import os
        
        # Use exact geometry provided by the layer.
        # Padding is now handled at the source (gemini_service) so preview and export are identical.
        left = int(layer.geometry.x * sx)
        top = int(layer.geometry.y * sy)
        width = int(layer.geometry.w * sx)
        height = int(layer.geometry.h * sy)
        
        original_image_path = layer.content.get('original_image')
        if original_image_path and os.path.exists(original_image_path):
            try:
                # Crop the image for the PPTX
                with Image.open(original_image_path) as img:
                    img_w, img_h = img.size
                    # Scale factor from 1280x720 space to actual image pixel grid
                    scale_w = img_w / 1280
                    scale_h = img_h / 720
                    
                    # Crop box using the source-padded geometry
                    crop_box = (
                        int(layer.geometry.x * scale_w),
                        int(layer.geometry.y * scale_h),
                        int((layer.geometry.x + layer.geometry.w) * scale_w),
                        int((layer.geometry.y + layer.geometry.h) * scale_h)
                    )
                    
                    # Ensure crop box is within image bounds
                    crop_box = (
                        max(0, crop_box[0]),
                        max(0, crop_box[1]),
                        min(img_w, crop_box[2]),
                        min(img_h, crop_box[3])
                    )
                    
                    if crop_box[2] > crop_box[0] and crop_box[3] > crop_box[1]:
                        cropped_img = img.crop(crop_box)
                        
                        # Save cropped image to a buffer
                        img_byte_arr = io.BytesIO()
                        cropped_img.save(img_byte_arr, format='PNG')
                        img_byte_arr.seek(0)
                        
                        slide.shapes.add_picture(img_byte_arr, left, top, width, height)
                        return
            except Exception as e:
                logger.error(f"Failed to crop and add image to PPTX: {e}")

        # Fallback to placeholder if cropping fails
        shape = slide.shapes.add_shape(1, left, top, width, height) # 1 = Rectangle
        shape.text = f"[Image: {layer.content.get('label', 'graphic')}]"

    def _add_shape(self, slide, layer, sx, sy):
        left = int(layer.geometry.x * sx)
        top = int(layer.geometry.y * sy)
        width = int(layer.geometry.w * sx)
        height = int(layer.geometry.h * sy)
        slide.shapes.add_shape(1, left, top, width, height)

pptx_service = PPTXService()
